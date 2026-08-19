"""PPT 演示文稿 8 大专业板式与 5 大主题渲染深度测试。"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from doc2mind.core.creator import (
    SlideLayoutType,
    extract_artifact,
)
from doc2mind.core.creator.exporters.pptx_exporter import PptxExporter


def test_pptx_parser_identifies_all_archetypes():
    marp_text = """
:::artifact type="pptx" title="全板式综合汇报" theme="emerald_green"
---
# DocMind 智能知识库
## 新一代离线优先创作平台
<!-- note: 封面开场词 -->
---
# 议程与核心目录
- 1. 架构革新
- 2. 核心性能看板
- 3. 三大演进支柱
- 4. 实施路线图
---
<!-- layout: cards -->
# 核心技术支柱
### 向量检索引擎
- 纯 CPU 轻量 ONNX
- 35MB 极小内存占用
### 智能图谱拓扑
- SQLite 原生存储
- 零额外网络开销
### 物理创作导出
- 原生 PPTX 几何渲染
- 100% 离线隐私安全
---
<!-- layout: metrics -->
# 性能指标看板
- 99.9% : 服务可用性与可靠度
- 10x : 向量混合检索加速比
- 12ms : 平均端到端检索耗时
- 0MB : GPU 显存依赖与开销
---
<!-- layout: timeline -->
# 阶段演进路线图
- 阶段一 : 知识库核心引擎研发
- 阶段二 : 创作工作台与多格式导出
- 阶段三 : 全自动 AI 整理与知识沉淀
---
<!-- layout: table -->
# 方案对比矩阵
| 维度 | DocMind | 传统方案 | 云端 SaaS |
| 显存需求 | 0MB | 8GB | 无法内网部署 |
| 导出质量 | 原生 PPTX 几何卡片 | 纯文本 | 依赖付费 API |
---
<!-- layout: quote -->
# 架构师箴言
> 极简、极致流畅与 100% 本地隐私安全，是个人智能知识库与创作工具的灵魂。
:::
"""
    artifact = extract_artifact(marp_text, default_type="pptx")
    assert artifact.title == "全板式综合汇报"
    assert artifact.theme == "emerald_green"
    assert len(artifact.slides) == 7

    assert artifact.slides[0].layout == SlideLayoutType.COVER
    assert artifact.slides[1].layout == SlideLayoutType.AGENDA
    assert artifact.slides[2].layout == SlideLayoutType.CARDS
    assert len(artifact.slides[2].cards) == 3
    assert artifact.slides[3].layout == SlideLayoutType.METRICS
    assert len(artifact.slides[3].metrics) == 4
    assert artifact.slides[4].layout == SlideLayoutType.TIMELINE
    assert len(artifact.slides[4].timeline_nodes) == 3
    assert artifact.slides[5].layout == SlideLayoutType.TABLE
    assert artifact.slides[6].layout == SlideLayoutType.QUOTE


@pytest.mark.parametrize("theme_name", ["tech_blue", "emerald_green", "modern_purple", "warm_orange", "dark_elegant"])
def test_pptx_exporter_renders_all_themes_and_layouts(tmp_path: Path, theme_name: str):
    marp_text = f"""
:::artifact type="pptx" title="主题测试: {theme_name}" theme="{theme_name}"
---
# 主题演示文稿
## 副标题说明
<!-- note: 备注测试 -->
---
<!-- layout: cards -->
# 多列卡片
### 卡片A
- 内容 1
- 内容 2
### 卡片B
- 内容 3
---
<!-- layout: metrics -->
# 核心指标
- 100% : 达标率
- 0ms : 延迟
---
<!-- layout: timeline -->
# 步骤
- Step 1 : 初始化
- Step 2 : 运行
"""
    artifact = extract_artifact(marp_text, default_type="pptx")
    out_file = tmp_path / f"test_{theme_name}.pptx"

    exporter = PptxExporter()
    res_path = exporter.export(artifact, out_file)

    assert res_path.exists()
    assert res_path.stat().st_size > 0

    # 验证生成的 PPTX 可被标准 python-pptx 打开且幻灯片页数一致
    prs = Presentation(str(res_path))
    assert len(prs.slides) == 4
    assert prs.slide_width.inches == pytest.approx(13.333, rel=1e-2)
    assert prs.slide_height.inches == pytest.approx(7.5, rel=1e-2)
    assert prs.slides[0].notes_slide.notes_text_frame.text == "备注测试"
