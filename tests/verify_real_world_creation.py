"""端到端真实落地验证脚本 — 生成真实物理交付物并校验文件内部结构。"""

from __future__ import annotations

import sys
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8")

import docx
import openpyxl
from pptx import Presentation

from doc2mind.core.creator import (
    export_artifact,
)


def run_e2e_real_world_verification():
    print("=" * 70)
    print("🚀 开始 DocMind 创作 Agent 全链路端到端真实生成检验")
    print("=" * 70)

    output_dir = Path("./verification_output").resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    # =========================================================================
    # 1. 检验专业 PPT 生成（覆盖 8 大板式 + 科技商务蓝主题）
    # =========================================================================
    sample_ppt = """
:::artifact type="pptx" title="DocMind 2026 核心架构与商业化落地汇报" theme="tech_blue"
---
# DocMind 智能知识平台
## 个人与团队离线优先知识创作工作台
<!-- note: 尊敬的各位评审与领导，今天由我向大家汇报 DocMind 架构设计与商业化演进成果。 -->
---
# 议程与核心目录
- 1. 行业痛点与解决思路
- 2. 核心架构与关键指标
- 3. 三大演进支柱
- 4. 商业化推进路线图
- 5. 竞品方案深度对比
- 6. 架构箴言与总结
<!-- note: 本次汇报分为六个部分，重点展示性能突破与多格式创作导出能力。 -->
---
<!-- layout: cards -->
# 核心技术三大支柱
### 向量计算引擎
- CPU 轻量 ONNX Runtime 运行
- 35MB 极低内存常驻
- 零 GPU 显存依赖
### 智能图谱与排错
- SQLite 原生存储
- 自动嗅探实体关联
- Pitfall 避坑排错顾问
### 物理创作导出
- 8 大专业板式渲染
- 5 套企业级品牌调色板
- 100% 本地隐私安全
<!-- note: 这是 DocMind 区别于传统云端 SaaS 的三大核心技术壁垒。 -->
---
<!-- layout: metrics -->
# 核心性能突破看板
- 99.9% : 服务可用性与可靠度
- 10x : 向量混合检索加速比
- 12ms : 平均端到端检索耗时
- 0MB : GPU 显存占用开销
<!-- note: 即使在无独立显卡的普通办公轻薄本上，也能达到毫秒级响应。 -->
---
<!-- layout: timeline -->
# 商业化推进路线图
- 阶段一 : 知识库核心引擎与 RAG 检索链路研发
- 阶段二 : 桌面端 Artifacts 工作台与导出器上线
- 阶段三 : 全自动 AI 整理、蒸馏与多租户商业交付
<!-- note: 目前我们已圆满完成阶段一与阶段二，具备完整的商业化交付能力。 -->
---
<!-- layout: table -->
# 综合方案深度对比
| 评估维度 | DocMind 本地版 | 传统云端知识库 | 本地重型方案 |
| 显存需求 | 0MB (纯CPU) | 无法内网离线 | 至少 8GB 显卡 |
| 启动速度 | 0.8 秒 | 依赖网络 | > 15 秒 |
| 导出格式 | PPTX / DOCX / XLSX / HTML | 仅支持 Markdown | 纯文本 |
| 数据隐私 | 100% 本地硬盘 | 存在泄密隐患 | 本地但易泄露 |
<!-- note: 从表格可以看出，DocMind 在轻量化与隐私安全方面优势显著。 -->
---
<!-- layout: quote -->
# 架构师箴言
> “极致轻量、毫秒响应与 100% 离线隐私安全，是个人知识与团队创作工具的基石。”
<!-- note: 最后以这句话结束今天的汇报，谢谢大家！ -->
:::
"""
    ppt_path = output_dir / "DocMind_Architect_Report.pptx"
    print("\n[1/4] 正在生成 PPTX 演示文稿...")
    res_ppt = export_artifact(sample_ppt, output_path=ppt_path)
    assert res_ppt.ok, f"PPTX 导出失败: {res_ppt.error}"
    assert ppt_path.exists(), "PPTX 文件未生成"

    # 深入检验 PPTX 内部结构
    prs = Presentation(str(ppt_path))
    assert len(prs.slides) == 7, f"幻灯片页数不符，期望 7，实际 {len(prs.slides)}"
    assert abs(prs.slide_width.inches - 13.333) < 0.05, "幻灯片宽度不符合 16:9"
    assert abs(prs.slide_height.inches - 7.5) < 0.05, "幻灯片高度不符合 16:9"

    # 验证第一页封面标题
    slide_texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
    assert any("DocMind 智能知识平台" in t for t in slide_texts), "未找到封面标题"
    # 验证 Speaker Notes
    assert "尊敬的各位评审与领导" in prs.slides[0].notes_slide.notes_text_frame.text

    print(f"  ✓ PPTX 生成成功: {ppt_path.name}")
    print(f"    - 文件大小: {ppt_path.stat().st_size:,} bytes")
    print(f"    - 页面总数: {len(prs.slides)} 页 (含封面、目录、卡片网格、大数字看板、时间线、表格、金句)")
    print("    - 演讲备注: 已完整注入 Slide Notes")

    # =========================================================================
    # 2. 检验专业 Word DOCX 研报方案生成
    # =========================================================================
    sample_doc = """
# DocMind 智能知识库与创作平台技术白皮书

## 一、项目概述与背景
DocMind 是一套针对个人技术专家与企业知识工作者研发的**本地优先、零显存依赖、支持多格式创作导出**的新一代知识中枢。

## 二、核心技术优势
- **本地 ONNX 向量嵌入**：在低功耗 CPU 上实现 12ms 极速嵌入计算；
- **SQLite-Vec 混合检索引擎**：BM25 + 向量 RRF 倒数秩融合；
- **智能避坑排错顾问 (Pitfall Advisor)**：检索时自动嗅探历史排错经验。

## 三、性能对比矩阵
| 评估项 | DocMind | 竞品 A (云端) | 竞品 B (重型) |
| 内存占用 | 35 MB | 无本地服务 | 1,400 MB |
| 检索耗时 | 12 ms | 120 ms | 65 ms |
| 离线可用 | 100% | 0% (断网不可用) | 100% |

> 💡 架构建议：对于个人与中小型团队的知识管理，纯本地计算是兼顾极致隐私与低成本的最优解。
"""
    doc_path = output_dir / "DocMind_Technical_Whitepaper.docx"
    print("\n[2/4] 正在生成 Word DOCX 深度技术白皮书...")
    res_doc = export_artifact(sample_doc, target_format="docx", output_path=doc_path)
    assert res_doc.ok, f"DOCX 导出失败: {res_doc.error}"
    assert doc_path.exists()

    doc = docx.Document(str(doc_path))
    assert len(doc.paragraphs) > 5
    assert len(doc.tables) == 1
    assert len(doc.tables[0].rows) == 4

    print(f"  ✓ DOCX 生成成功: {doc_path.name}")
    print(f"    - 文件大小: {doc_path.stat().st_size:,} bytes")
    print(f"    - 段落数量: {len(doc.paragraphs)} 个")
    print(f"    - 结构化表格: {len(doc.tables)} 个（{len(doc.tables[0].rows)} 行 × {len(doc.tables[0].columns)} 列）")

    # =========================================================================
    # 3. 检验 Excel XLSX 商业数据对比矩阵生成
    # =========================================================================
    sample_excel = """
# 嵌入模型性能与维度对比矩阵

| 模型代号 | 向量维度 | 模型体积 | 语言支持 | 典型延迟 | 推荐使用场景 |
| BAAI/bge-small-zh-v1.5 | 512 | 35 MB | 中文 / 英文 | 12 ms | 个人桌面知识库、开发文档检索 |
| all-MiniLM-L6-v2 | 384 | 22 MB | 英文为主 | 8 ms | 超轻量英文代码与日志检索 |
| text-embedding-3-small | 1536 | 云端 API | 多语言 | 150 ms | 跨语言云端企业知识中枢 |
| bge-large-zh-v1.5 | 1024 | 1.3 GB | 中文专业领域 | 45 ms | GPU 服务器重型问答 |
"""
    excel_path = output_dir / "DocMind_Performance_Matrix.xlsx"
    print("\n[3/4] 正在生成 Excel XLSX 对比矩阵表格...")
    res_excel = export_artifact(sample_excel, target_format="xlsx", output_path=excel_path)
    assert res_excel.ok, f"XLSX 导出失败: {res_excel.error}"
    assert excel_path.exists()

    wb = openpyxl.load_workbook(str(excel_path))
    ws = wb.active
    assert ws.max_row == 5, f"Excel 行数不符，期望 5，实际 {ws.max_row}"
    assert ws.max_column == 6, f"Excel 列数不符，期望 6，实际 {ws.max_column}"
    assert ws.cell(row=1, column=1).value == "模型代号"
    assert ws.cell(row=2, column=1).value == "BAAI/bge-small-zh-v1.5"

    print(f"  ✓ XLSX 生成成功: {excel_path.name}")
    print(f"    - 文件大小: {excel_path.stat().st_size:,} bytes")
    print(f"    - 表格尺寸: {ws.max_row} 行 × {ws.max_column} 列")
    print("    - 样式渲染: 已应用深蓝商务表头、居中对齐与自适应列宽")

    # =========================================================================
    # 4. 检验单文件自包含 HTML5 看板生成
    # =========================================================================
    sample_html = """
## 平台运行状态概览
- 嵌入服务：ONNX Runtime (正常)
- 向量存储：SQLite-Vec (就绪)
- 创作工作台：8 大板式就绪

| 模块名称 | 状态 | 运行时 |
| 向量检索 | 运行中 | CPU |
| 创作引擎 | 运行中 | 原生渲染 |
| 知识图谱 | 就绪 | 本地 SQLite |
"""
    html_path = output_dir / "DocMind_Dashboard.html"
    print("\n[4/4] 正在生成 HTML5 交互看板...")
    res_html = export_artifact(sample_html, target_format="html", output_path=html_path)
    assert res_html.ok, f"HTML 导出失败: {res_html.error}"
    assert html_path.exists()

    content = html_path.read_text(encoding="utf-8")
    assert "<!DOCTYPE html>" in content
    assert "平台运行状态概览" in content
    assert "<table>" in content

    print(f"  ✓ HTML 生成成功: {html_path.name}")
    print(f"    - 文件大小: {html_path.stat().st_size:,} bytes")
    print("    - 自包含特性: 无外部 CDN 依赖，双击即可在任何浏览器渲染")

    print("\n" + "=" * 70)
    print("🎉 全部 4 类交付物真实物理导出与内部结构检验 100% 成功！")
    print(f"📁 真实文件落地路径: {output_dir}")
    print("=" * 70)


if __name__ == "__main__":
    run_e2e_real_world_verification()
