"""创作导出引擎单元测试。"""

from __future__ import annotations

from pathlib import Path

import docx
import openpyxl
from pptx import Presentation

from doc2mind.core.creator import (
    ArtifactType,
    export_artifact,
    extract_artifact,
)
from doc2mind.core.creator.exporters.docx_exporter import DocxExporter
from doc2mind.core.creator.exporters.excel_exporter import ExcelExporter
from doc2mind.core.creator.exporters.html_exporter import HtmlExporter
from doc2mind.core.creator.exporters.pptx_exporter import PptxExporter


def test_extract_artifact_pptx_and_slides():
    sample_marp = """
:::artifact type="pptx" title="DocMind 架构解析"
---
# DocMind 架构解析
## 本地智能知识库与创作平台
<!-- note: 各位好，今天为大家汇报 DocMind 架构 -->
---
# 核心特性
- 本地 ONNX 向量嵌入
- SQLite-Vec 混合检索
- 零依赖极速启动
<!-- note: 重点讲解轻量化与本地隐私优势 -->
---
# 性能对比
| 维度 | DocMind | 传统方案 |
| 内存占用 | 35MB | 1.2GB |
| 启动时间 | 0.8s | 15s |
:::
"""
    artifact = extract_artifact(sample_marp, default_type="pptx")
    assert artifact.artifact_type == ArtifactType.PPTX
    assert artifact.title == "DocMind 架构解析"
    assert len(artifact.slides) == 3
    assert artifact.slides[0].is_cover is True
    assert artifact.slides[0].subtitle == "本地智能知识库与创作平台"
    assert "各位好" in artifact.slides[0].speaker_notes
    assert len(artifact.slides[1].bullet_points) == 3
    assert artifact.slides[2].table_data is not None
    assert len(artifact.slides[2].table_data) == 3


def test_pptx_exporter_creates_valid_file(tmp_path: Path):
    sample_marp = """
# DocMind 架构汇报
## 本地轻量知识平台
<!-- note: 封面开场白 -->
---
# 架构优势
- 100% 本地隐私安全
- CPU 极轻量嵌入
<!-- note: 第二页演讲重点 -->
"""
    artifact = extract_artifact(sample_marp, default_type="pptx")
    out_file = tmp_path / "test_presentation.pptx"

    exporter = PptxExporter()
    res_path = exporter.export(artifact, out_file)

    assert res_path.exists()
    assert res_path.stat().st_size > 0

    # 用 python-pptx 重新打开验证
    prs = Presentation(str(res_path))
    assert len(prs.slides) == 2
    assert prs.slides[0].notes_slide.notes_text_frame.text == "封面开场白"


def test_docx_exporter_creates_valid_file(tmp_path: Path):
    sample_doc = """
# 智能知识库技术研报

## 1. 概述与背景
DocMind 是一套本地优先的轻量级个人知识库与创作平台。

## 2. 方案对比
| 指标 | DocMind | 竞品 A |
| 嵌入耗时 | 12ms | 85ms |
| 显存占用 | 0MB | 2000MB |

> 💡 架构建议：对于桌面端个人工具，CPU 嵌入足以达到极致流畅体验。
"""
    artifact = extract_artifact(sample_doc, default_type="docx")
    out_file = tmp_path / "test_report.docx"

    exporter = DocxExporter()
    res_path = exporter.export(artifact, out_file)

    assert res_path.exists()
    assert res_path.stat().st_size > 0

    # 用 python-docx 重新打开验证
    doc = docx.Document(str(res_path))
    assert len(doc.paragraphs) > 3
    assert len(doc.tables) == 1
    assert len(doc.tables[0].rows) == 3


def test_excel_exporter_creates_valid_file(tmp_path: Path):
    sample_table = """
# 嵌入模型性能对比矩阵

| 模型名称 | 向量维度 | 模型体积 | 适用语言 | 推荐场景 |
| BAAI/bge-small-zh-v1.5 | 512 | 35MB | 中文/英文 | 个人桌面端与本地知识库 |
| all-MiniLM-L6-v2 | 384 | 22MB | 英文 | 极低资源英文检索 |
| text-embedding-3-small | 1536 | 云端 API | 多语言 | 云端高精度场景 |
"""
    artifact = extract_artifact(sample_table, default_type="xlsx")
    out_file = tmp_path / "test_matrix.xlsx"

    exporter = ExcelExporter()
    res_path = exporter.export(artifact, out_file)

    assert res_path.exists()
    assert res_path.stat().st_size > 0

    # 用 openpyxl 重新打开验证
    wb = openpyxl.load_workbook(str(res_path))
    ws = wb.active
    assert ws.max_row == 4
    assert ws.max_column == 5
    assert ws.cell(row=1, column=1).value == "模型名称"
    assert ws.cell(row=2, column=1).value == "BAAI/bge-small-zh-v1.5"


def test_html_exporter_creates_valid_file(tmp_path: Path):
    sample_web = """
## 核心特性总览
- 本地 0 显存占用
- 混合 RRF 检索

| 功能项 | 状态 |
| 向量库 | 已就绪 |
| 创作引擎 | 已就绪 |
"""
    artifact = extract_artifact(sample_web, default_type="html")
    out_file = tmp_path / "test_dashboard.html"

    exporter = HtmlExporter()
    res_path = exporter.export(artifact, out_file)

    assert res_path.exists()
    html_text = res_path.read_text(encoding="utf-8")
    assert "<!DOCTYPE html>" in html_text
    assert "核心特性总览" in html_text
    assert "<table>" in html_text


def test_html_slideshow_exporter_all_themes(tmp_path: Path):
    sample_marp = """
:::artifact type="pptx" title="HTML 放映测试" theme="dark_elegant"
---
# 封面页
## 交互式 Web 放映
<!-- note: 演讲者小抄 -->
---
# 核心指标
- 99.9%: 准确率
- 15ms: 延迟
---
# 团队寄语
> 技术点亮生活，极致创造价值。
:::
"""
    artifact = extract_artifact(sample_marp, default_type="pptx")
    out_file = tmp_path / "test_slideshow.html"

    exporter = HtmlExporter()
    res_path = exporter.export(artifact, out_file)

    assert res_path.exists()
    html_text = res_path.read_text(encoding="utf-8")
    assert "DocMind SlideShow" in html_text
    assert "演讲者实时提词小抄" in html_text
    assert "99.9%" in html_text
    assert "技术点亮生活" in html_text


def test_export_artifact_factory_auto_dispatch(tmp_path: Path):
    content = """
:::artifact type="pptx" title="自动化导出测试"
---
# 自动化测试 PPT
- 测试项 1
- 测试项 2
:::
"""
    target = tmp_path / "auto_ppt.pptx"
    result = export_artifact(content=content, output_path=target)

    assert result.ok is True
    assert result.artifact_type == "pptx"
    assert Path(result.file_path).exists()
    assert result.file_size_bytes > 0
