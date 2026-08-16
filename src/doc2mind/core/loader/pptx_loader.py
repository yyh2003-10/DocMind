"""PowerPoint 加载器 — 基于 `python-pptx`。

特点：
- 逐 slide 遍历，输出 H1 幻灯片标题
- 遍历 shape：`has_text_frame` → 文本框；`has_table` → 表格行
- 标题检测启发式：字号 > slide 平均 × 1.5 或 shape 顶部 20% 区域
- 备注页 (notes_slide) 文本一并提取
- 表格转换为 markdown table_row 序列

局限性：
- SmartArt / 图表内容无法提取
- 图片中文字需用 image_loader 的 OCR
"""

from __future__ import annotations

import hashlib
from pathlib import Path

from doc2mind.core.loader.base import Loader, LoaderError, make_source
from doc2mind.core.models import (
    DocFormat,
    DocumentElement,
    ElementType,
    LoadedDocument,
)


def _shape_text(shape) -> str:  # noqa: ANN001 — pptx shape 类型动态
    """提取 shape 中的纯文本，保留段落换行。"""
    if not shape.has_text_frame:
        return ""
    lines: list[str] = []
    for para in shape.text_frame.paragraphs:
        runs = [r.text for r in para.runs if r.text]
        line = "".join(runs) if runs else (para.text or "")
        if line:
            lines.append(line)
    return "\n".join(lines).strip()


def _shape_avg_font_size(shape) -> float:  # noqa: ANN001
    """计算 shape 内所有 run 的平均字号，无显式字号返回 0。"""
    if not shape.has_text_frame:
        return 0.0
    sizes: list[float] = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            font = run.font
            if font.size is not None:
                sizes.append(float(font.size.pt))
    if not sizes:
        return 0.0
    return sum(sizes) / len(sizes)


class PptxLoader(Loader):
    """PowerPoint 文档加载器（python-pptx 实现）。"""

    supported_extensions = ("pptx", "ppt")

    def extract(self, path: Path) -> LoadedDocument:
        try:
            from pptx import Presentation
        except ImportError as e:
            raise LoaderError(
                "python-pptx 未安装。请运行：pip install python-pptx"
            ) from e

        if not path.exists():
            raise LoaderError(f"文件不存在: {path}")

        try:
            data = path.read_bytes()
            file_hash = hashlib.md5(data).hexdigest()
            prs = Presentation(path)
            elements: list[DocumentElement] = []

            for slide_idx, slide in enumerate(prs.slides, start=1):
                # 幻灯片级标题
                elements.append(
                    DocumentElement(
                        content=f"# Slide {slide_idx}",
                        type=ElementType.HEADING,
                        metadata={
                            "type": "heading",
                            "level": 1,
                            "slide": slide_idx,
                            "source_format": DocFormat.PPTX.value,
                        },
                    )
                )

                # 收集所有 shape 的字号用于启发式
                shape_sizes: list[float] = []
                shapes = list(slide.shapes)
                for shape in shapes:
                    sz = _shape_avg_font_size(shape)
                    if sz > 0:
                        shape_sizes.append(sz)
                slide_avg_size = (
                    sum(shape_sizes) / len(shape_sizes) if shape_sizes else 18.0
                )

                # slide 顶部 20% 区域启发式（EMU 单位）
                try:
                    slide_height = prs.slide_height
                except Exception:  # noqa: BLE001
                    slide_height = 0
                top_threshold = slide_height * 0.2 if slide_height else 0

                for shape in shapes:
                    # 表格
                    if shape.has_table:
                        table = shape.table
                        for r_idx, row in enumerate(table.rows):
                            cells = [
                                _cell_to_table_text(cell) for cell in row.cells
                            ]
                            if not any(c.strip() for c in cells):
                                continue
                            elements.append(
                                DocumentElement(
                                    content="| " + " | ".join(cells) + " |",
                                    type=ElementType.TABLE_ROW,
                                    metadata={
                                        "type": "table_row",
                                        "slide": slide_idx,
                                        "row_index": r_idx + 1,
                                        "cols": len(cells),
                                        "source_format": DocFormat.PPTX.value,
                                    },
                                )
                            )
                        continue

                    # 文本框
                    text = _shape_text(shape)
                    if not text:
                        continue

                    avg_size = _shape_avg_font_size(shape)
                    # 标题启发式：字号大 OR 位于顶部 20%
                    is_title = (
                        avg_size > slide_avg_size * 1.5 and avg_size > 20
                    ) or (
                        top_threshold > 0
                        and getattr(shape, "top", None) is not None
                        and shape.top <= top_threshold
                    )

                    if is_title:
                        elements.append(
                            DocumentElement(
                                content=text,
                                type=ElementType.HEADING,
                                metadata={
                                    "type": "heading",
                                    "level": 2,
                                    "slide": slide_idx,
                                    "font_size": round(avg_size, 2),
                                    "source_format": DocFormat.PPTX.value,
                                },
                            )
                        )
                    else:
                        elements.append(
                            DocumentElement(
                                content=text,
                                type=ElementType.PARAGRAPH,
                                metadata={
                                    "type": "paragraph",
                                    "slide": slide_idx,
                                    "font_size": round(avg_size, 2)
                                    if avg_size
                                    else None,
                                    "source_format": DocFormat.PPTX.value,
                                },
                            )
                        )

                # 备注页文本
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        elements.append(
                            DocumentElement(
                                content=notes_text,
                                type=ElementType.PARAGRAPH,
                                metadata={
                                    "type": "paragraph",
                                    "slide": slide_idx,
                                    "role": "speaker_notes",
                                    "source_format": DocFormat.PPTX.value,
                                },
                            )
                        )

            return LoadedDocument(
                source=make_source(path),
                format=DocFormat.PPTX,
                elements=elements,
                page_count=len(prs.slides) if prs.slides else None,
                size_bytes=len(data),
                file_hash=file_hash,
            )
        except LoaderError:
            raise
        except Exception as e:  # noqa: BLE001
            raise LoaderError(f"PowerPoint 解析失败 ({path.name}): {e}") from e


def _cell_to_table_text(cell) -> str:  # noqa: ANN001
    """提取表格单元格文本（合并所有段落）。"""
    try:
        text = cell.text_frame.text
    except Exception:  # noqa: BLE001
        return ""
    return (text or "").replace("\n", " ").strip()
