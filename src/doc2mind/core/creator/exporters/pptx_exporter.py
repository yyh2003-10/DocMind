"""PPTX 原生幻灯片导出引擎 V2 — 基于 python-pptx。

特性：
- 默认 16:9 现代超清宽屏（13.333 × 7.5 英寸）；
- 5 大企业级主题配色方案（科技蓝、自然绿、智能紫、活力橙、极简暗黑）；
- 8 大专业板式几何排版渲染（Cover, Agenda, Cards, Metrics, Timeline, Table, Quote, General）；
- 智能自适应网格、圆角卡片容器、色块装饰条与徽章；
- 100% 完整注入 Slide Speaker Notes（演讲备注）。
"""

from __future__ import annotations

import logging
from pathlib import Path

from doc2mind.core.creator.models import (
    ArtifactModel,
    MetricItem,
    SlideCardItem,
    SlideLayoutType,
    SlideModel,
    TimelineNodeItem,
)
from doc2mind.core.creator.parser import parse_pptx_slides
from doc2mind.core.creator.themes import PptTheme, RgbColor, get_theme

logger = logging.getLogger("doc2mind.creator.pptx")


def _to_pptx_rgb(c: RgbColor):
    """转为 python-pptx RGBColor。"""
    from pptx.dml.color import RGBColor
    return RGBColor(c.r, c.g, c.b)


class PptxExporter:
    """专业级 PPTX 幻灯片几何排版渲染引擎。"""

    def __init__(self) -> None:
        pass

    def export(self, artifact: ArtifactModel, output_path: Path) -> Path:
        """编译生成高保真 .pptx 物理文件。"""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        slides = artifact.slides
        if not slides:
            slides = parse_pptx_slides(artifact.raw_content)

        if not slides:
            slides = [
                SlideModel(
                    index=1,
                    title=artifact.title or "DocMind 演示文稿",
                    subtitle="基于 DocMind 智能知识库生成",
                    layout=SlideLayoutType.COVER,
                    is_cover=True,
                )
            ]

        # 获取选定的企业主题
        theme = get_theme(artifact.theme)

        for slide_model in slides:
            slide = prs.slides.add_slide(blank_layout)

            # 1. 绘制页面全屏背景底色
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = _to_pptx_rgb(theme.bg)
            bg_shape.line.fill.background()

            # 2. 写入演讲备注 Speaker Notes
            if slide_model.speaker_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_model.speaker_notes

            # 3. 按板式类型分流渲染
            layout = slide_model.layout
            if slide_model.is_cover or layout == SlideLayoutType.COVER:
                self._render_cover(slide, slide_model, theme, len(slides))
            elif layout == SlideLayoutType.AGENDA:
                self._render_agenda(slide, slide_model, theme, len(slides))
            elif layout == SlideLayoutType.CARDS:
                self._render_cards(slide, slide_model, theme, len(slides))
            elif layout == SlideLayoutType.METRICS:
                self._render_metrics(slide, slide_model, theme, len(slides))
            elif layout == SlideLayoutType.TIMELINE:
                self._render_timeline(slide, slide_model, theme, len(slides))
            elif layout == SlideLayoutType.TABLE:
                self._render_table(slide, slide_model, theme, len(slides))
            elif layout == SlideLayoutType.QUOTE:
                self._render_quote(slide, slide_model, theme, len(slides))
            else:
                self._render_general(slide, slide_model, theme, len(slides))

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        logger.info("PPTX 导出成功: %s, 共 %d 页, 主题: %s", output_path, len(slides), theme.name)
        return output_path

    # =========================================================================
    # 板式 1: 科技商务封面 (Cover)
    # =========================================================================
    def _render_cover(self, slide, model: SlideModel, theme: PptTheme, total_slides: int) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        # 顶部/左侧品牌几何装饰带
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.2)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = _to_pptx_rgb(theme.accent)
        accent_bar.line.fill.background()

        # 封面背景卡片 (大号优雅浮层)
        cover_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.3)
        )
        cover_card.fill.solid()
        cover_card.fill.fore_color.rgb = _to_pptx_rgb(theme.card_bg)
        cover_card.line.color.rgb = _to_pptx_rgb(theme.card_border)
        cover_card.line.width = Pt(1.5)

        # 左侧装饰竖条
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(1.8), Inches(0.18), Inches(3.8)
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = _to_pptx_rgb(theme.primary)
        stripe.line.fill.background()

        # 主标题
        title_box = slide.shapes.add_textbox(Inches(1.8), Inches(1.8), Inches(10.0), Inches(2.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = model.title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = _to_pptx_rgb(theme.primary)
        p.font.name = "Microsoft YaHei"

        # 副标题
        sub_box = slide.shapes.add_textbox(Inches(1.8), Inches(3.8), Inches(10.0), Inches(1.0))
        stf = sub_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = model.subtitle or "DocMind 智能知识库深度创作出品"
        sp.font.size = Pt(20)
        sp.font.color.rgb = _to_pptx_rgb(theme.text_muted)
        sp.font.name = "Microsoft YaHei"

        # 底部作者 / 机构信息
        footer_box = slide.shapes.add_textbox(Inches(1.8), Inches(5.4), Inches(10.0), Inches(0.5))
        fp = footer_box.text_frame.paragraphs[0]
        fp.text = "💡 DocMind Studio  |  依托知识库与全网深度研报系统"
        fp.font.size = Pt(12)
        fp.font.color.rgb = _to_pptx_rgb(theme.secondary)
        fp.font.name = "Microsoft YaHei"

    # =========================================================================
    # 板式 2: 序号矩阵目录 (Agenda)
    # =========================================================================
    def _render_agenda(self, slide, model: SlideModel, theme: PptTheme, total_slides: int) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        self._draw_header(slide, model, theme, total_slides)

        items = model.bullet_points if model.bullet_points else ["概述与背景", "核心技术架构", "关键特性与成果", "总结与实施规划"]
        count = len(items)
        cols = 2 if count > 3 else 1
        rows = (count + cols - 1) // cols

        start_x = Inches(1.2)
        start_y = Inches(2.0)
        card_w = Inches(5.3) if cols == 2 else Inches(10.9)
        card_h = Inches(min(1.2, 4.5 / max(1, rows)))
        gap_x = Inches(0.5)
        gap_y = Inches(0.3)

        for i, item in enumerate(items):
            r = i // cols
            c = i % cols
            x = start_x + c * (card_w + gap_x)
            y = start_y + r * (card_h + gap_y)

            # 目录条目卡片
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = _to_pptx_rgb(theme.card_bg)
            card.line.color.rgb = _to_pptx_rgb(theme.card_border)

            # 左侧序号胶囊
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.2), y + (card_h - Inches(0.6)) / 2, Inches(0.8), Inches(0.6))
            badge.fill.solid()
            badge.fill.fore_color.rgb = _to_pptx_rgb(theme.primary)
            badge.line.fill.background()
            bp = badge.text_frame.paragraphs[0]
            bp.text = f"{i+1:02d}"
            bp.font.size = Pt(16)
            bp.font.bold = True
            bp.font.color.rgb = _to_pptx_rgb(RgbColor(255, 255, 255))
            bp.font.name = "Arial"

            # 目录标题
            t_box = slide.shapes.add_textbox(x + Inches(1.2), y + (card_h - Inches(0.6)) / 2, card_w - Inches(1.4), Inches(0.6))
            tp = t_box.text_frame.paragraphs[0]
            tp.text = item
            tp.font.size = Pt(16)
            tp.font.bold = True
            tp.font.color.rgb = _to_pptx_rgb(theme.text_main)
            tp.font.name = "Microsoft YaHei"

    # =========================================================================
    # 板式 3: 多列卡片网格 (Cards 2/3/4-Column)
    # =========================================================================
    def _render_cards(self, slide, model: SlideModel, theme: PptTheme, total_slides: int) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        self._draw_header(slide, model, theme, total_slides)

        cards = model.cards
        if not cards and model.bullet_points:
            cards = [SlideCardItem(title=f"模块 0{i+1}", content=b) for i, b in enumerate(model.bullet_points[:4])]

        count = min(max(len(cards), 1), 4)
        total_width = Inches(11.733)
        start_x = Inches(0.8)
        start_y = Inches(2.0)
        card_h = Inches(4.7)

        gap = Inches(0.3)
        card_w = (total_width - gap * (count - 1)) / count

        for i, card_item in enumerate(cards[:count]):
            x = start_x + i * (card_w + gap)

            # 主卡片背景
            card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, card_w, card_h)
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = _to_pptx_rgb(theme.card_bg)
            card_shape.line.color.rgb = _to_pptx_rgb(theme.card_border)
            card_shape.line.width = Pt(1.5)

            # 顶部强调色细条
            top_stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y, card_w, Inches(0.12))
            top_stripe.fill.solid()
            top_stripe.fill.fore_color.rgb = _to_pptx_rgb(theme.primary if i % 2 == 0 else theme.secondary)
            top_stripe.line.fill.background()

            # 序号微标
            badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), start_y + Inches(0.3), Inches(0.45), Inches(0.45))
            badge.fill.solid()
            badge.fill.fore_color.rgb = _to_pptx_rgb(theme.primary)
            badge.line.fill.background()
            bp = badge.text_frame.paragraphs[0]
            bp.text = str(i + 1)
            bp.font.size = Pt(12)
            bp.font.bold = True
            bp.font.color.rgb = _to_pptx_rgb(RgbColor(255, 255, 255))
            bp.font.name = "Arial"

            # 卡片标题
            t_box = slide.shapes.add_textbox(x + Inches(0.75), start_y + Inches(0.25), card_w - Inches(0.9), Inches(0.6))
            tf = t_box.text_frame
            tf.word_wrap = True
            tp = tf.paragraphs[0]
            tp.text = card_item.title
            tp.font.size = Pt(16)
            tp.font.bold = True
            tp.font.color.rgb = _to_pptx_rgb(theme.text_main)
            tp.font.name = "Microsoft YaHei"

            # 卡片正文/列表
            body_box = slide.shapes.add_textbox(x + Inches(0.2), start_y + Inches(1.0), card_w - Inches(0.4), card_h - Inches(1.2))
            btf = body_box.text_frame
            btf.word_wrap = True

            if card_item.bullets:
                for b_idx, b in enumerate(card_item.bullets):
                    p = btf.paragraphs[0] if b_idx == 0 else btf.add_paragraph()
                    p.text = f"• {b}"
                    p.font.size = Pt(13)
                    p.font.color.rgb = _to_pptx_rgb(theme.text_main)
                    p.font.name = "Microsoft YaHei"
                    p.space_after = Pt(8)
            elif card_item.content:
                p = btf.paragraphs[0]
                p.text = card_item.content
                p.font.size = Pt(13)
                p.font.color.rgb = _to_pptx_rgb(theme.text_main)
                p.font.name = "Microsoft YaHei"

    # =========================================================================
    # 板式 4: 大数字 KPI / 核心数据看板 (Metrics)
    # =========================================================================
    def _render_metrics(self, slide, model: SlideModel, theme: PptTheme, total_slides: int) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        self._draw_header(slide, model, theme, total_slides)

        metrics = model.metrics
        if not metrics:
            metrics = [MetricItem(value="100%", label="核心达标率", description="生产环境验证稳定通过")]

        count = min(max(len(metrics), 1), 4)
        total_width = Inches(11.733)
        start_x = Inches(0.8)
        start_y = Inches(2.2)
        card_h = Inches(4.3)

        gap = Inches(0.35)
        card_w = (total_width - gap * (count - 1)) / count

        for i, m in enumerate(metrics[:count]):
            x = start_x + i * (card_w + gap)

            card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, card_w, card_h)
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = _to_pptx_rgb(theme.card_bg)
            card_shape.line.color.rgb = _to_pptx_rgb(theme.card_border)
            card_shape.line.width = Pt(1.5)

            # 超大号核心数值 (38~44pt Bold)
            num_box = slide.shapes.add_textbox(x + Inches(0.1), start_y + Inches(0.6), card_w - Inches(0.2), Inches(1.2))
            np = num_box.text_frame.paragraphs[0]
            np.alignment = PP_ALIGN.CENTER
            np.text = m.value
            np.font.size = Pt(40)
            np.font.bold = True
            np.font.color.rgb = _to_pptx_rgb(theme.primary)
            np.font.name = "Arial"

            # 指标名称
            label_box = slide.shapes.add_textbox(x + Inches(0.1), start_y + Inches(1.9), card_w - Inches(0.2), Inches(0.6))
            lp = label_box.text_frame.paragraphs[0]
            lp.alignment = PP_ALIGN.CENTER
            lp.text = m.label
            lp.font.size = Pt(17)
            lp.font.bold = True
            lp.font.color.rgb = _to_pptx_rgb(theme.text_main)
            lp.font.name = "Microsoft YaHei"

            # 指标释义
            if m.description or i < len(model.bullet_points):
                desc_text = m.description or model.bullet_points[i]
                desc_box = slide.shapes.add_textbox(x + Inches(0.2), start_y + Inches(2.6), card_w - Inches(0.4), Inches(1.3))
                dp = desc_box.text_frame.paragraphs[0]
                dp.alignment = PP_ALIGN.CENTER
                dp.text = desc_text
                dp.font.size = Pt(12)
                dp.font.color.rgb = _to_pptx_rgb(theme.text_muted)
                dp.font.name = "Microsoft YaHei"

    # =========================================================================
    # 板式 5: 横向时间线 / 阶段路线图 (Timeline)
    # =========================================================================
    def _render_timeline(self, slide, model: SlideModel, theme: PptTheme, total_slides: int) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        self._draw_header(slide, model, theme, total_slides)

        nodes = model.timeline_nodes
        if not nodes:
            nodes = [TimelineNodeItem(stage=f"Step 0{i+1}", title=b) for i, b in enumerate(model.bullet_points[:4])]

        count = min(max(len(nodes), 1), 4)
        total_width = Inches(11.733)
        start_x = Inches(0.8)
        axis_y = Inches(2.8)
        gap = Inches(0.3)
        node_w = (total_width - gap * (count - 1)) / count

        # 绘制主时间轴连接线
        axis_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, start_x + Inches(0.5), axis_y + Inches(0.2), total_width - Inches(1.0), Inches(0.06)
        )
        axis_line.fill.solid()
        axis_line.fill.fore_color.rgb = _to_pptx_rgb(theme.secondary)
        axis_line.line.fill.background()

        for i, node in enumerate(nodes[:count]):
            x = start_x + i * (node_w + gap)

            # 阶段药丸徽章
            pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + (node_w - Inches(1.5)) / 2, axis_y, Inches(1.5), Inches(0.45))
            pill.fill.solid()
            pill.fill.fore_color.rgb = _to_pptx_rgb(theme.primary)
            pill.line.fill.background()
            pp = pill.text_frame.paragraphs[0]
            pp.alignment = PP_ALIGN.CENTER
            pp.text = node.stage
            pp.font.size = Pt(13)
            pp.font.bold = True
            pp.font.color.rgb = _to_pptx_rgb(RgbColor(255, 255, 255))
            pp.font.name = "Microsoft YaHei"

            # 阶段卡片容器
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, axis_y + Inches(0.7), node_w, Inches(3.2))
            card.fill.solid()
            card.fill.fore_color.rgb = _to_pptx_rgb(theme.card_bg)
            card.line.color.rgb = _to_pptx_rgb(theme.card_border)

            # 阶段标题
            t_box = slide.shapes.add_textbox(x + Inches(0.15), axis_y + Inches(0.85), node_w - Inches(0.3), Inches(0.6))
            tp = t_box.text_frame.paragraphs[0]
            tp.alignment = PP_ALIGN.CENTER
            tp.text = node.title
            tp.font.size = Pt(15)
            tp.font.bold = True
            tp.font.color.rgb = _to_pptx_rgb(theme.text_main)
            tp.font.name = "Microsoft YaHei"

            # 阶段详情要点
            if node.details:
                d_box = slide.shapes.add_textbox(x + Inches(0.15), axis_y + Inches(1.5), node_w - Inches(0.3), Inches(2.2))
                dtf = d_box.text_frame
                dtf.word_wrap = True
                for b_idx, b in enumerate(node.details):
                    p = dtf.paragraphs[0] if b_idx == 0 else dtf.add_paragraph()
                    p.text = f"• {b}"
                    p.font.size = Pt(12)
                    p.font.color.rgb = _to_pptx_rgb(theme.text_main)
                    p.font.name = "Microsoft YaHei"

    # =========================================================================
    # 板式 6: 商务对比矩阵 (Table)
    # =========================================================================
    def _render_table(self, slide, model: SlideModel, theme: PptTheme, total_slides: int) -> None:
        from pptx.util import Inches, Pt

        self._draw_header(slide, model, theme, total_slides)

        t_data = model.table_data or [["对比项", "DocMind", "传统方案"], ["检索耗时", "12ms", "120ms"], ["本地隐私", "100% 离线", "依赖公网"]]
        rows = len(t_data)
        cols = max(len(r) for r in t_data)

        t_left = Inches(1.0)
        t_top = Inches(2.0)
        t_width = Inches(11.333)
        t_height = Inches(min(4.5, rows * 0.7))

        table_shape = slide.shapes.add_table(rows, cols, t_left, t_top, t_width, t_height)
        table = table_shape.table

        for r_idx, row in enumerate(t_data):
            for c_idx, cell_value in enumerate(row):
                if c_idx < cols:
                    cell = table.cell(r_idx, c_idx)
                    cell.text = cell_value

                    # 表头行深色填充，交替行浅色斑马纹
                    if r_idx == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _to_pptx_rgb(theme.primary)
                    elif r_idx % 2 == 1:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _to_pptx_rgb(theme.card_bg)
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _to_pptx_rgb(theme.bg)

                    for cp in cell.text_frame.paragraphs:
                        cp.font.name = "Microsoft YaHei"
                        if r_idx == 0:
                            cp.font.size = Pt(14)
                            cp.font.bold = True
                            cp.font.color.rgb = _to_pptx_rgb(RgbColor(255, 255, 255))
                        else:
                            cp.font.size = Pt(13)
                            cp.font.color.rgb = _to_pptx_rgb(theme.text_main)

    # =========================================================================
    # 板式 7: 金句强调卡片 (Quote)
    # =========================================================================
    def _render_quote(self, slide, model: SlideModel, theme: PptTheme, total_slides: int) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        self._draw_header(slide, model, theme, total_slides)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(10.333), Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = _to_pptx_rgb(theme.card_bg)
        card.line.color.rgb = _to_pptx_rgb(theme.card_border)
        card.line.width = Pt(1.5)

        # 左侧强调粗线
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.2), Inches(0.2), Inches(4.2))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = _to_pptx_rgb(theme.accent)
        stripe.line.fill.background()

        q_box = slide.shapes.add_textbox(Inches(2.2), Inches(2.8), Inches(9.0), Inches(2.8))
        qtf = q_box.text_frame
        qtf.word_wrap = True
        qp = qtf.paragraphs[0]
        qp.text = f"“ {model.quote_text or model.bullet_points[0] if model.bullet_points else '核心结论与洞察'} ”"
        qp.font.size = Pt(24)
        qp.font.bold = True
        qp.font.color.rgb = _to_pptx_rgb(theme.text_main)
        qp.font.name = "Microsoft YaHei"

    # =========================================================================
    # 板式 8: 经典图文正文 (General)
    # =========================================================================
    def _render_general(self, slide, model: SlideModel, theme: PptTheme, total_slides: int) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        self._draw_header(slide, model, theme, total_slides)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = _to_pptx_rgb(theme.card_bg)
        card.line.color.rgb = _to_pptx_rgb(theme.card_border)
        card.line.width = Pt(1.5)

        content_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(10.9), Inches(4.2))
        ctf = content_box.text_frame
        ctf.word_wrap = True

        for b_idx, bullet in enumerate(model.bullet_points):
            bp = ctf.paragraphs[0] if b_idx == 0 else ctf.add_paragraph()
            bp.text = f"•   {bullet}"
            bp.font.size = Pt(17)
            bp.font.name = "Microsoft YaHei"
            bp.font.color.rgb = _to_pptx_rgb(theme.text_main)
            bp.space_after = Pt(14)

    def _draw_header(self, slide, model: SlideModel, theme: PptTheme, total_slides: int) -> None:
        """公共顶部标题栏与页码。"""
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        # 标题左侧指示色标
        indicator = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.6), Inches(0.12), Inches(0.8))
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = _to_pptx_rgb(theme.primary)
        indicator.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.5), Inches(11.0), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = model.title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = _to_pptx_rgb(theme.primary)
        p.font.name = "Microsoft YaHei"

        # 底部页码
        page_box = slide.shapes.add_textbox(Inches(11.0), Inches(6.8), Inches(1.5), Inches(0.4))
        pp = page_box.text_frame.paragraphs[0]
        pp.alignment = PP_ALIGN.RIGHT
        pp.text = f"{model.index} / {total_slides}"
        pp.font.size = Pt(11)
        pp.font.color.rgb = _to_pptx_rgb(theme.text_muted)
        pp.font.name = "Arial"
